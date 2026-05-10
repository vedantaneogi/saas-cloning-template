"""Generate the dummy attachment fixtures shipped with the default seed.

Run from repo root with the api venv:
    apps/api/.venv/bin/python apps/api/seeds/samples/_generate.py

Produces (per senior's preview-format list — pdf / image / pptx / docs / excel + csv):
    Q4-Sales-Report.xlsx   — multi-sheet workbook with formulas & formatting
    Project-Brief.pdf      — single-page PDF summary
    Logo.png               — small inline image
    Q1-Pipeline.csv        — comma-separated pipeline data
    Status-Memo.docx       — Word memo with bullets / numbering
    Q4-Review.pptx         — PowerPoint deck

These files are loaded by `seed_loader.py` into the in-memory attachment
store so anyone running the default seed can preview real binary content
without uploading their own files first.
"""
from __future__ import annotations

import os
from pathlib import Path

OUT = Path(__file__).resolve().parent


def make_xlsx() -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill

    wb = Workbook()
    ws = wb.active
    ws.title = "Q4 Summary"

    header_fill = PatternFill("solid", fgColor="0078D4")
    header_font = Font(bold=True, color="FFFFFF", size=11)
    money = '"$"#,##0'

    ws.append(["Region", "Product", "Units Sold", "Unit Price", "Revenue"])
    for cell in ws[1]:
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center", vertical="center")

    rows = [
        ("North America", "Outlook Pro", 1240, 49.99),
        ("North America", "Outlook Lite", 3120, 9.99),
        ("EMEA",          "Outlook Pro",  890, 49.99),
        ("EMEA",          "Outlook Lite", 2410, 9.99),
        ("APAC",          "Outlook Pro",  670, 49.99),
        ("APAC",          "Outlook Lite", 1980, 9.99),
        ("LATAM",         "Outlook Pro",  310, 49.99),
        ("LATAM",         "Outlook Lite", 720,  9.99),
    ]
    for r in rows:
        ws.append([*r, None])

    for row_idx in range(2, 2 + len(rows)):
        ws.cell(row=row_idx, column=4).number_format = money
        ws.cell(row=row_idx, column=5).value = f"=C{row_idx}*D{row_idx}"
        ws.cell(row=row_idx, column=5).number_format = money

    total_row = 2 + len(rows)
    ws.cell(row=total_row, column=1).value = "Total"
    ws.cell(row=total_row, column=1).font = Font(bold=True)
    ws.cell(row=total_row, column=3).value = f"=SUM(C2:C{total_row - 1})"
    ws.cell(row=total_row, column=5).value = f"=SUM(E2:E{total_row - 1})"
    ws.cell(row=total_row, column=5).number_format = money
    ws.cell(row=total_row, column=5).font = Font(bold=True)

    for col_letter, width in zip("ABCDE", (16, 16, 12, 12, 14)):
        ws.column_dimensions[col_letter].width = width

    # second sheet — pipeline
    ws2 = wb.create_sheet("Pipeline")
    ws2.append(["Account", "Stage", "ARR", "Owner"])
    for cell in ws2[1]:
        cell.fill = header_fill
        cell.font = header_font
    pipeline = [
        ("Acme Corp",        "Negotiation", 120000, "Frank Miller"),
        ("Globex",           "Closed Won",   86000, "Carol Williams"),
        ("Initech",          "Discovery",    42000, "Bob Smith"),
        ("Soylent",          "Proposal",    175000, "Alice Wilson"),
        ("Massive Dynamic",  "Negotiation", 240000, "Frank Miller"),
    ]
    for r in pipeline:
        ws2.append(r)
    for row_idx in range(2, 2 + len(pipeline)):
        ws2.cell(row=row_idx, column=3).number_format = money
    for col_letter, width in zip("ABCD", (20, 16, 14, 18)):
        ws2.column_dimensions[col_letter].width = width

    out = OUT / "Q4-Sales-Report.xlsx"
    wb.save(out)
    print(f"wrote {out} ({out.stat().st_size} bytes)")


def make_pdf() -> None:
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import (Paragraph, SimpleDocTemplate, Spacer, Table,
                                    TableStyle)
    from reportlab.lib import colors

    out = OUT / "Project-Brief.pdf"
    doc = SimpleDocTemplate(str(out), pagesize=LETTER,
                            leftMargin=54, rightMargin=54, topMargin=54, bottomMargin=54)
    styles = getSampleStyleSheet()
    story = []
    story.append(Paragraph("Q4 Outlook Migration — Project Brief", styles["Title"]))
    story.append(Spacer(1, 12))
    story.append(Paragraph("<b>Owner:</b> Frank Miller &nbsp;&nbsp; <b>Status:</b> In progress &nbsp;&nbsp; <b>Target:</b> Dec 15", styles["BodyText"]))
    story.append(Spacer(1, 8))
    story.append(Paragraph(
        "This document summarises the scope, milestones, and risk profile for the "
        "Outlook clone migration. It is intended as a single-page reference for "
        "stakeholders who need to confirm the rollout plan ahead of the December "
        "freeze.", styles["BodyText"]))
    story.append(Spacer(1, 12))
    story.append(Paragraph("Milestones", styles["Heading2"]))
    data = [
        ["#", "Milestone",            "Owner",         "Due"],
        ["1", "Mail triage parity",   "Frank Miller",  "Nov 14"],
        ["2", "Calendar recurrence",  "Carol Williams","Nov 21"],
        ["3", "Universal search",     "Bob Smith",     "Nov 28"],
        ["4", "Attachment previews",  "Alice Wilson",  "Dec 05"],
        ["5", "Customer pilot",       "Frank Miller",  "Dec 12"],
    ]
    t = Table(data, colWidths=[24, 220, 110, 90])
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0078D4")),
        ("TEXTCOLOR",  (0, 0), (-1, 0), colors.white),
        ("FONTNAME",   (0, 0), (-1, 0), "Helvetica-Bold"),
        ("ALIGN",      (0, 0), (0, -1), "CENTER"),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.whitesmoke, colors.white]),
        ("GRID",       (0, 0), (-1, -1), 0.4, colors.HexColor("#EDEBE9")),
        ("FONTSIZE",   (0, 0), (-1, -1), 10),
    ]))
    story.append(t)
    story.append(Spacer(1, 12))
    story.append(Paragraph("Risks", styles["Heading2"]))
    story.append(Paragraph(
        "1. Recurrence rule edge cases (BYDAY, exception dates).<br/>"
        "2. TipTap version drift — extensions vs starter-kit.<br/>"
        "3. Time-zone handling for cross-region invitees.", styles["BodyText"]))
    doc.build(story)
    print(f"wrote {out} ({out.stat().st_size} bytes)")


def make_png() -> None:
    from PIL import Image, ImageDraw, ImageFont

    img = Image.new("RGB", (480, 240), (15, 108, 189))
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 36)
    except OSError:
        font = ImageFont.load_default()
    draw.text((28, 90), "Outlook Clone", fill="white", font=font)
    try:
        sub = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 18)
    except OSError:
        sub = ImageFont.load_default()
    draw.text((28, 140), "Q4 marketing banner", fill=(220, 230, 245), font=sub)
    out = OUT / "Logo.png"
    img.save(out, format="PNG")
    print(f"wrote {out} ({out.stat().st_size} bytes)")


def make_csv() -> None:
    out = OUT / "Q1-Pipeline.csv"
    out.write_text(
        "Account,Stage,ARR,Owner,Close Date\n"
        "Acme Corp,Negotiation,120000,Frank Miller,2026-03-15\n"
        "Globex,Closed Won,86000,Carol Williams,2026-02-21\n"
        "Initech,Discovery,42000,Bob Smith,2026-04-30\n"
        "Soylent,Proposal,175000,Alice Wilson,2026-04-10\n"
        "Massive Dynamic,Negotiation,240000,Frank Miller,2026-03-28\n"
        "Hooli,Closed Won,98000,Carol Williams,2026-02-08\n",
        encoding="utf-8",
    )
    print(f"wrote {out} ({out.stat().st_size} bytes)")


def make_docx() -> None:
    from docx import Document
    from docx.shared import Pt, RGBColor

    doc = Document()
    title = doc.add_heading("Q4 Outlook Migration — Status Memo", level=0)
    for run in title.runs:
        run.font.color.rgb = RGBColor(0x00, 0x78, 0xD4)

    doc.add_paragraph(
        "Owner: Frank Miller   ·   Status: In progress   ·   Target: Dec 15"
    )
    doc.add_paragraph(
        "This memo captures the rollout state for the Outlook clone migration "
        "ahead of the December freeze. Send replies on this thread to confirm "
        "or push back on the milestones below."
    )

    doc.add_heading("Highlights this week", level=2)
    for line in (
        "Recurring events ship-ready — series delete + scope edit verified.",
        "Search universal dropdown live in production; Files tab now opens previews.",
        "Boomerang follow-up wired end-to-end; reply suppression confirmed.",
    ):
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(line).font.size = Pt(11)

    doc.add_heading("Open risks", level=2)
    for line in (
        "Cross-region timezone handling for invitees.",
        "TipTap version drift between starter-kit and font-size extension.",
        "DLP false-positive rate on long internal threads.",
    ):
        p = doc.add_paragraph(style="List Number")
        p.add_run(line).font.size = Pt(11)

    out = OUT / "Status-Memo.docx"
    doc.save(out)
    print(f"wrote {out} ({out.stat().st_size} bytes)")


def make_pptx() -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]

    # Slide 1 — title
    s1 = prs.slides.add_slide(title_layout)
    s1.shapes.title.text = "Outlook Clone — Q4 Review"
    s1.placeholders[1].text = "Frank Miller · Acme Corp · 2026-05"

    # Slide 2 — agenda
    s2 = prs.slides.add_slide(bullet_layout)
    s2.shapes.title.text = "Agenda"
    body = s2.placeholders[1].text_frame
    body.text = "Recap of P0 milestones"
    for line in ("Calendar recurrence + series delete",
                 "Universal search dropdown",
                 "Real attachment previews (xlsx, pdf, docx, pptx)",
                 "RSVP & scheduling assistant"):
        p = body.add_paragraph()
        p.text = line

    # Slide 3 — metrics
    s3 = prs.slides.add_slide(bullet_layout)
    s3.shapes.title.text = "Metrics this quarter"
    body3 = s3.placeholders[1].text_frame
    body3.text = "Coverage of CSV P0/P1 features"
    for line in ("P0 coverage: 100%",
                 "P1 coverage: 70% → 90% (this batch)",
                 "Open: PowerPoint preview, sensitivity badge"):
        p = body3.add_paragraph()
        p.text = line

    # Slide 4 — next steps
    s4 = prs.slides.add_slide(bullet_layout)
    s4.shapes.title.text = "Next steps"
    body4 = s4.placeholders[1].text_frame
    body4.text = "Phase 4 candidates"
    for line in ("Sensitivity labels in reading pane",
                 "Add-ins / Quick Steps polish",
                 "Public release branch"):
        p = body4.add_paragraph()
        p.text = line

    out = OUT / "Q4-Review.pptx"
    prs.save(out)
    print(f"wrote {out} ({out.stat().st_size} bytes)")


if __name__ == "__main__":
    make_xlsx()
    make_pdf()
    make_png()
    make_csv()
    make_docx()
    make_pptx()
